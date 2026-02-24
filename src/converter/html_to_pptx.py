"""HTML to PPTX converter using html2image for accurate rendering."""

import os
import re
import tempfile
from io import BytesIO
from typing import Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image
import logging

logger = logging.getLogger(__name__)

# Try to import html2image, fallback to None if not available
try:
    from html2image import Html2Image
    HTI_AVAILABLE = True
except ImportError:
    HTI_AVAILABLE = False
    logger.warning("html2image not available. Install with: pip install html2image")


def render_html_to_png(
    html: str,
    output_path: str,
    width: int = 1280,
    height: int = 720,
) -> str:
    """Render an HTML document to a PNG image.

    This is used by the MCP preview tool to let models validate slide layout
    before exporting to PPTX.
    """
    if not HTI_AVAILABLE:
        raise RuntimeError("html2image is not installed")

    abs_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(abs_path), exist_ok=True)

    # Ensure HTML has proper body styles to fill viewport
    html = _ensure_fullpage_styles(html, width, height)

    hti = Html2Image(size=(width, height), output_path=os.path.dirname(abs_path))

    filename = os.path.basename(abs_path)
    if not filename.lower().endswith(".png"):
        filename = f"{filename}.png"
        abs_path = os.path.join(os.path.dirname(abs_path), filename)

    paths = hti.screenshot(
        html_str=html,
        save_as=filename,
        size=(width, height),
    )

    if not paths:
        raise RuntimeError("Failed to render HTML to PNG")

    return abs_path


def _ensure_fullpage_styles(html: str, width: int, height: int) -> str:
    """Ensure HTML fills the entire viewport without borders."""
    # Add critical styles if not present
    base_style = f"""
    <style>
        html, body {{
            margin: 0 !important;
            padding: 0 !important;
            width: {width}px !important;
            height: {height}px !important;
            overflow: hidden !important;
        }}
        * {{
            box-sizing: border-box;
        }}
    </style>
    """
    
    # If no <head>, add one with styles
    if "<head>" not in html.lower():
        html = html.replace("<html", f"<html><head>{base_style}</head>", 1)
    elif "</head>" in html:
        # Insert before </head>
        html = html.replace("</head>", f"{base_style}</head>", 1)
    else:
        # Fallback: prepend to body
        html = html.replace("<body", f"{base_style}<body", 1)
    
    return html


class HTMLToPPTXConverter:
    """
    Converts HTML slides to PowerPoint presentation.
    
    Uses html2image to render HTML to images first, then adds them as slides.
    This ensures accurate rendering of CSS, Tailwind, and complex layouts.
    """
    
    def __init__(
        self,
        width: int = 1280,
        height: int = 720,
        fit_mode: str = "contain",
        background_color: str = "#FFFFFF",
        render_bleed_px: int = 0,
    ):
        """
        Initialize converter with slide dimensions.
        
        Args:
            width: Slide width in pixels
            height: Slide height in pixels
        """
        self.width = width
        self.height = height
        self.fit_mode = (fit_mode or "contain").lower()
        self.background_color = background_color
        self.render_bleed_px = max(0, int(render_bleed_px or 0))
        
        # Convert pixels to EMUs (English Metric Units)
        # 1 inch = 914400 EMUs, 96 DPI
        self.width_emu = int(width * 914400 / 96)
        self.height_emu = int(height * 914400 / 96)
        
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = self.width_emu
        self.prs.slide_height = self.height_emu
        
        # Initialize html2image
        self.hti = None
        if HTI_AVAILABLE:
            try:
                self.hti = Html2Image(
                    size=(width, height),
                    output_path=tempfile.gettempdir()
                )
            except Exception as e:
                logger.warning(f"Failed to initialize html2image: {e}")
    
    def _render_html_to_image(self, html_content: str, output_path: Optional[str] = None) -> Optional[bytes]:
        """Render HTML content to an image."""
        if not self.hti:
            logger.error("html2image not available")
            return None

        try:
            html_content = _ensure_fullpage_styles(html_content, self.width, self.height)

            import uuid
            filename = f"slide_{uuid.uuid4().hex[:8]}"

            bleed = self.render_bleed_px
            render_size = (self.width + 2 * bleed, self.height + 2 * bleed)

            paths = self.hti.screenshot(
                html_str=html_content,
                save_as=f"{filename}.png",
                size=render_size,
            )

            if not paths:
                return None

            image_path = paths[0]

            with open(image_path, "rb") as f:
                image_bytes = f.read()

            if bleed > 0:
                try:
                    with Image.open(BytesIO(image_bytes)) as img:
                        img = img.convert("RGBA")
                        # Center-crop to target slide size
                        left = max(0, (img.width - self.width) // 2)
                        top = max(0, (img.height - self.height) // 2)
                        img = img.crop((left, top, left + self.width, top + self.height))
                        out = BytesIO()
                        img.save(out, format="PNG")
                        image_bytes = out.getvalue()
                except Exception as e:
                    logger.warning(f"Bleed crop failed: {e}")

            image_bytes = self._normalize_rendered_image(image_bytes)

            if output_path:
                import shutil
                with open(output_path, "wb") as out:
                    out.write(image_bytes)

            try:
                os.remove(image_path)
            except Exception:
                pass

            return image_bytes

        except Exception as e:
            logger.error(f"Failed to render HTML to image: {e}")
            return None

    def _normalize_rendered_image(self, image_bytes: bytes) -> bytes:
        """Trim borders and fit to slide size.

        Production behavior:
        - Trim constant-color borders (often white/black bars from rendering).
        - Fit the rendered content into the slide using `fit_mode`:
          - contain: never crop (letterbox/pillarbox using theme background)
          - cover: fill slide and crop overflow (best for full-bleed visuals)
          - stretch: fill slide without preserving aspect (can distort)
        """
        try:
            with Image.open(BytesIO(image_bytes)) as img:
                img = img.convert("RGBA")

                img = self._trim_uniform_borders(img)

                target_w, target_h = self.width, self.height
                mode = self.fit_mode

                if mode == "stretch":
                    if img.size != (target_w, target_h):
                        img = img.resize((target_w, target_h), resample=Image.LANCZOS)

                elif mode == "cover":
                    src_ratio = img.width / img.height
                    tgt_ratio = target_w / target_h
                    if src_ratio > tgt_ratio:
                        new_w = int(img.height * tgt_ratio)
                        left = (img.width - new_w) // 2
                        img = img.crop((left, 0, left + new_w, img.height))
                    elif src_ratio < tgt_ratio:
                        new_h = int(img.width / tgt_ratio)
                        top = (img.height - new_h) // 2
                        img = img.crop((0, top, img.width, top + new_h))

                    if img.size != (target_w, target_h):
                        img = img.resize((target_w, target_h), resample=Image.LANCZOS)

                else:
                    # contain (default): never crop; letterbox with background
                    src_ratio = img.width / img.height
                    tgt_ratio = target_w / target_h

                    if src_ratio >= tgt_ratio:
                        new_w = target_w
                        new_h = max(1, int(target_w / src_ratio))
                    else:
                        new_h = target_h
                        new_w = max(1, int(target_h * src_ratio))

                    if img.size != (new_w, new_h):
                        img = img.resize((new_w, new_h), resample=Image.LANCZOS)

                    bg = Image.new("RGBA", (target_w, target_h), _parse_hex_color(self.background_color) + (255,))
                    x = (target_w - new_w) // 2
                    y = (target_h - new_h) // 2
                    bg.paste(img, (x, y), img)
                    img = bg

                out = BytesIO()
                img.save(out, format="PNG")
                return out.getvalue()
        except Exception as e:
            logger.warning(f"Image normalization failed: {e}")
            return image_bytes

    def _trim_uniform_borders(self, img: Image.Image) -> Image.Image:
        """Trim constant-color borders from an RGBA image."""
        try:
            w, h = img.size
            if w < 10 or h < 10:
                return img

            px = img.load()

            def close(c1, c2, tol: int = 8) -> bool:
                return all(abs(int(c1[i]) - int(c2[i])) <= tol for i in range(3))

            corners = [px[0, 0], px[w - 1, 0], px[0, h - 1], px[w - 1, h - 1]]
            border = max(set(corners), key=corners.count)

            def row_is_border(y: int) -> bool:
                step = max(1, w // 32)
                for x in range(0, w, step):
                    if not close(px[x, y], border):
                        return False
                return True

            def col_is_border(x: int) -> bool:
                step = max(1, h // 32)
                for y in range(0, h, step):
                    if not close(px[x, y], border):
                        return False
                return True

            top = 0
            while top < h - 1 and row_is_border(top):
                top += 1

            bottom = h - 1
            while bottom > top and row_is_border(bottom):
                bottom -= 1

            left = 0
            while left < w - 1 and col_is_border(left):
                left += 1

            right = w - 1
            while right > left and col_is_border(right):
                right -= 1

            if top > 0 or left > 0 or bottom < h - 1 or right < w - 1:
                pad = 2
                left = max(0, left - pad)
                top = max(0, top - pad)
                right = min(w - 1, right + pad)
                bottom = min(h - 1, bottom + pad)
                img = img.crop((left, top, right + 1, bottom + 1))

            return img
        except Exception as e:
            logger.warning(f"Uniform border trim failed: {e}")
            return img

    def add_slide_from_html(self, html_content: str, background_color: str = "#FFFFFF") -> bool:
        """Add a slide to the presentation from HTML content."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        image_bytes = self._render_html_to_image(html_content)
        if image_bytes:
            return self._add_fullslide_image(slide, image_bytes)

        return self._set_background_color(slide, background_color)

    def _add_fullslide_image(self, slide, image_bytes: bytes) -> bool:
        """Add an image that fills the entire slide."""
        try:
            image_stream = BytesIO(image_bytes)
            left = Emu(0)
            top = Emu(0)
            width = Emu(self.width_emu)
            height = Emu(self.height_emu)
            slide.shapes.add_picture(image_stream, left, top, width, height)
            return True
        except Exception as e:
            logger.error(f"Failed to add image to slide: {e}")
            return False

    def _set_background_color(self, slide, color: str) -> bool:
        """Set slide background color."""
        try:
            if not color:
                return False
            c = color.strip()
            if c.startswith('#'):
                hex_color = c[1:]
                if len(hex_color) == 3:
                    hex_color = ''.join([ch * 2 for ch in hex_color])
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)

                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)
                return True
        except Exception as e:
            logger.error(f"Failed to set background color: {e}")
        return False

    def save(self, output_path: str) -> str:
        """Save the presentation to a file and return absolute path."""
        abs_path = os.path.abspath(output_path)
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        self.prs.save(abs_path)
        return abs_path

    def get_bytes(self) -> bytes:
        """Get the presentation as bytes."""
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


def _parse_hex_color(color: str) -> tuple[int, int, int]:
    c = (color or "#FFFFFF").strip()
    if not c.startswith("#"):
        c = f"#{c}"
    hex_color = c[1:]
    if len(hex_color) == 3:
        hex_color = "".join([ch * 2 for ch in hex_color])
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return r, g, b
    except Exception:
        return 255, 255, 255

    


def convert_html_slides_to_pptx(
    slides: list[dict],
    output_path: str,
    width: int = 1280,
    height: int = 720,
    theme_config: Optional[dict] = None,
    fit_mode: str = "contain",
    safe_margin_px: int = 0,
    render_bleed_px: int = 0,
) -> str:
    """
    Convert a list of HTML slides to a PowerPoint presentation.
    
    Args:
        slides: List of slide dictionaries with 'html_content' and 'index' keys
        output_path: Path to save the .pptx file
        width: Slide width in pixels
        height: Slide height in pixels
        theme_config: Theme configuration with background color (fallback)
    
    Returns:
        Path to the generated presentation
    """
    # Get fallback background color from theme
    bg_color = "#FFFFFF"
    if theme_config and 'background' in theme_config:
        bg_color = theme_config['background']

    converter = HTMLToPPTXConverter(
        width,
        height,
        fit_mode=fit_mode,
        background_color=bg_color,
        render_bleed_px=render_bleed_px,
    )
    
    # Sort slides by index and add them
    sorted_slides = sorted(slides, key=lambda s: s.get('index', 0))
    
    for slide_data in sorted_slides:
        html_content = slide_data.get('html_content', '')
        if html_content:
            if safe_margin_px and safe_margin_px > 0:
                html_content = _inject_safe_margin(html_content, safe_margin_px)
            converter.add_slide_from_html(html_content, bg_color)
    
    return converter.save(output_path)


def _inject_safe_margin(html: str, margin_px: int) -> str:
    # Inject safe padding WITHOUT breaking <body ...> attributes.
    style = f"""<style>
      .mcp-safe {{ padding: {int(margin_px)}px !important; width: 100%; height: 100%; box-sizing: border-box; }}
    </style>"""

    # Add style into <head> if present.
    if re.search(r"</head>", html, flags=re.IGNORECASE):
        html = re.sub(r"</head>", style + "\n</head>", html, count=1, flags=re.IGNORECASE)
    else:
        html = style + "\n" + html

    # Wrap body inner content.
    m = re.search(r"<body\b[^>]*>", html, flags=re.IGNORECASE)
    if not m:
        return html

    insert_at = m.end()
    html = html[:insert_at] + "\n<div class=\"mcp-safe\">\n" + html[insert_at:]
    html = re.sub(r"</body>", "\n</div>\n</body>", html, count=1, flags=re.IGNORECASE)
    return html


# Alternative converter using Playwright directly (more reliable)
class PlaywrightConverter:
    """
    Alternative converter using Playwright directly for more control.
    """
    
    def __init__(self, width: int = 1280, height: int = 720):
        self.width = width
        self.height = height
        self.width_emu = int(width * 914400 / 96)
        self.height_emu = int(height * 914400 / 96)
        
        self.prs = Presentation()
        self.prs.slide_width = self.width_emu
        self.prs.slide_height = self.height_emu
    
    async def render_html(self, html_content: str) -> Optional[bytes]:
        """Render HTML using Playwright."""
        try:
            from playwright.async_api import async_playwright
            
            async with async_playwright() as p:
                browser = await p.chromium.launch()
                page = await browser.new_page(
                    viewport={'width': self.width, 'height': self.height}
                )
                
                await page.set_content(html_content)
                await page.wait_for_load_state('networkidle')
                
                # Take screenshot
                screenshot = await page.screenshot(
                    type='png',
                    full_page=False
                )
                
                await browser.close()
                return screenshot
                
        except Exception as e:
            logger.error(f"Playwright rendering failed: {e}")
            return None
    
    def add_slide_from_image(self, image_bytes: bytes) -> bool:
        """Add a slide from image bytes."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        try:
            image_stream = BytesIO(image_bytes)
            left = Inches(0)
            top = Inches(0)
            width = Inches(self.width / 96)
            height = Inches(self.height / 96)
            
            slide.shapes.add_picture(image_stream, left, top, width, height)
            return True
        except Exception as e:
            logger.error(f"Failed to add slide: {e}")
            return False
    
    def save(self, output_path: str) -> str:
        """Save presentation."""
        abs_path = os.path.abspath(output_path)
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        self.prs.save(abs_path)
        return abs_path
